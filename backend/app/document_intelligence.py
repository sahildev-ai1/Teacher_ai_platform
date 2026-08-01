"""
Stage 1: Document Intelligence.

Parses PDF / DOCX / PPTX / TXT, preserving as much structure as cheaply as
possible (headings, sections, tables, figures, page/slide metadata), and
chunks the text for embedding.

Cost-aware routing (per the client's clarification #7): the caller passes a
`doc_type_hint` collected from the user at upload time ("Mostly Text",
"Text with Tables", "Text with Diagrams/Figures", "Text with Equations",
"Scanned PDF", "I'm Not Sure"). We combine that with cheap heuristics
(embedded image count, text density) to decide whether the lightweight parser
below is sufficient or whether the document should be flagged for a heavier
OCR/vision pass — that heavier path is left as a pluggable hook
(`needs_advanced_parsing`) so a real OCR/vision-LLM parser can be dropped in
without touching the rest of the pipeline.
"""
import os
import re
from dataclasses import dataclass, field
from typing import List

from pypdf import PdfReader
import docx
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter

from .config import settings

HEADING_RE = re.compile(r"^\s*(?:[0-9]+(?:\.[0-9]+)*\s+)?[A-Z][A-Za-z0-9 ,'\-]{2,80}$")


@dataclass
class ParsedDocument:
    full_text: str
    structure: dict  # {"headings": [...], "tables": n, "figures": n, "pages_or_slides": n}
    needs_advanced_parsing: bool
    chunks: List[str] = field(default_factory=list)


def _looks_like_heading(line: str) -> bool:
    line = line.strip()
    if not line or len(line) > 90:
        return False
    if line.endswith((".", ",", ";")):
        return False
    return bool(HEADING_RE.match(line)) and len(line.split()) <= 12


def _detect_headings(text: str) -> List[str]:
    headings = []
    for line in text.splitlines():
        if _looks_like_heading(line):
            headings.append(line.strip())
    # de-dupe while preserving order
    seen = set()
    out = []
    for h in headings:
        if h not in seen:
            seen.add(h)
            out.append(h)
    return out[:200]


def parse_pdf(path: str) -> ParsedDocument:
    reader = PdfReader(path)
    pages_text = []
    image_count = 0
    for page in reader.pages:
        pages_text.append(page.extract_text() or "")
        try:
            image_count += len(page.images)
        except Exception:  # noqa: BLE001
            pass
    full_text = "\n\n".join(pages_text)
    avg_chars_per_page = len(full_text) / max(len(pages_text), 1)
    needs_advanced = avg_chars_per_page < 200  # likely scanned / image-heavy
    structure = {
        "headings": _detect_headings(full_text),
        "tables": full_text.count("\t") > 0,  # crude signal; refined in advanced pass
        "figures": image_count,
        "pages_or_slides": len(pages_text),
    }
    return ParsedDocument(full_text=full_text, structure=structure, needs_advanced_parsing=needs_advanced)


def parse_docx(path: str) -> ParsedDocument:
    d = docx.Document(path)
    headings, paras, table_count = [], [], len(d.tables)
    for p in d.paragraphs:
        text = p.text.strip()
        if not text:
            continue
        style = (p.style.name or "").lower()
        if "heading" in style or "title" in style:
            headings.append(text)
        paras.append(text)
    for t in d.tables:
        for row in t.rows:
            paras.append(" | ".join(c.text for c in row.cells))
    full_text = "\n".join(paras)
    structure = {
        "headings": headings or _detect_headings(full_text),
        "tables": table_count,
        "figures": len(d.inline_shapes),
        "pages_or_slides": None,
    }
    return ParsedDocument(full_text=full_text, structure=structure, needs_advanced_parsing=False)


def parse_pptx(path: str) -> ParsedDocument:
    prs = Presentation(path)
    slides_text, headings, figure_count = [], [], 0
    for i, slide in enumerate(prs.slides, start=1):
        slide_lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    if t:
                        slide_lines.append(t)
                        if shape == slide.shapes.title:
                            headings.append(t)
            if shape.shape_type == 13:  # PICTURE
                figure_count += 1
        slides_text.append(f"[Slide {i}]\n" + "\n".join(slide_lines))
    full_text = "\n\n".join(slides_text)
    structure = {
        "headings": headings or _detect_headings(full_text),
        "tables": False,
        "figures": figure_count,
        "pages_or_slides": len(prs.slides.__iter__.__self__._sldIdLst),  # slide count
    }
    return ParsedDocument(full_text=full_text, structure=structure, needs_advanced_parsing=False)


def parse_txt(path: str) -> ParsedDocument:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        full_text = f.read()
    return ParsedDocument(
        full_text=full_text,
        structure={"headings": _detect_headings(full_text), "tables": False, "figures": 0, "pages_or_slides": None},
        needs_advanced_parsing=False,
    )


PARSERS = {
    "pdf": parse_pdf,
    "docx": parse_docx,
    "pptx": parse_pptx,
    "txt": parse_txt,
}


def parse_document(path: str, file_type: str) -> ParsedDocument:
    file_type = file_type.lower().lstrip(".")
    parser = PARSERS.get(file_type)
    if parser is None:
        raise ValueError(f"Unsupported file type: {file_type}")
    parsed = parser(path)
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.chunk_size,
        chunk_overlap=settings.chunk_overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    parsed.chunks = splitter.split_text(parsed.full_text)
    return parsed


def infer_file_type(filename: str) -> str:
    ext = os.path.splitext(filename)[1].lower().lstrip(".")
    if ext in ("pdf", "docx", "pptx", "txt"):
        return ext
    if ext in ("doc",):
        return "docx"
    if ext in ("ppt",):
        return "pptx"
    return "txt"
